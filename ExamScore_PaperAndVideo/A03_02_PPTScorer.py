import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.shapes.picture import Picture
import cv2
import numpy as np
import win32com.client as win32


def convert_ppt_to_pptx(root_dir):
    """
    将 .ppt 文件转换为 .pptx 文件，并删除原始 .ppt 文件
    """
    powerpoint = win32.Dispatch("PowerPoint.Application")
    powerpoint.Visible = True  # 改为 True，避免异常

    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.lower().endswith('.ppt') and not filename.lower().endswith('.pptx'):
                ppt_path = os.path.join(dirpath, filename)
                pptx_path = os.path.join(dirpath, os.path.splitext(filename)[0] + '.pptx')

                print(f"正在转换: {ppt_path} 到 {pptx_path}")

                try:
                    # 打开 .ppt 文件
                    presentation = powerpoint.Presentations.Open(ppt_path)
                    # 转换为 .pptx
                    presentation.SaveAs(pptx_path, 24)  # 24 表示 pptx 格式
                    presentation.Close()

                    # 删除原始 .ppt 文件
                    os.remove(ppt_path)
                    print(f"已删除原始文件: {ppt_path}")

                except Exception as e:
                    print(f"处理文件时发生错误: {ppt_path}\n错误信息: {e}")

    # 退出 PowerPoint 应用程序
    powerpoint.Quit()


def is_default_office_template(prs: Presentation) -> bool:
    """
    判断是否使用了默认 Office 模板
    """
    for master in prs.slide_masters:
        if master.name and "Office Theme" in master.name:
            return True
    return False


from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
import json
import os


def is_section_header(slide):
    """
    判断是否为 Section Header 页
    """
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                text = shape.text.strip()
                if text and len(text.split()) <= 10:
                    return True
    return False


def extract_slide_content(slide):
    """
    提取单页内容
    """
    contents = []
    figures = 0
    tables = 0

    for shape in slide.shapes:

        # 文本
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                text = p.text.strip()
                if text:
                    contents.append({
                        "type": "text",
                        "level": p.level,
                        "text": text
                    })

        # 图片
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            figures += 1

        # 表格
        if shape.has_table:
            tables += 1

    # Notes
    notes = None
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()

    return contents, figures, tables, notes


def ppt_to_paper_json(ppt_path, output_json):
    prs = Presentation(ppt_path)

    paper = {
        "meta": {
            "source": os.path.basename(ppt_path),
            "slide_count": len(prs.slides)
        },
        "sections": []
    }

    current_section = None
    section_id = 0

    for idx, slide in enumerate(prs.slides, start=1):

        # 识别 Section Header
        if is_section_header(slide):
            section_id += 1

            title = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                    title = shape.text.strip()

            current_section = {
                "section_id": section_id,
                "title": title if title else f"Section {section_id}",
                "slides": [idx],
                "content": [],
                "figures": 0,
                "tables": 0,
                "notes": ""
            }

            paper["sections"].append(current_section)
            continue

        # 普通内容页
        if current_section is None:
            # 没有 section header，自动生成 Abstract / Introduction
            section_id += 1
            current_section = {
                "section_id": section_id,
                "title": "Introduction",
                "slides": [],
                "content": [],
                "figures": 0,
                "tables": 0,
                "notes": ""
            }
            paper["sections"].append(current_section)

        contents, figs, tabs, notes = extract_slide_content(slide)

        current_section["slides"].append(idx)
        current_section["content"].extend(contents)
        current_section["figures"] += figs
        current_section["tables"] += tabs

        if notes:
            current_section["notes"] += notes + "\n"

    with open(output_json, "w", encoding="utf-8") as f:
        json.dump(paper, f, ensure_ascii=False, indent=2)

    print(f"✅ JSON saved to {output_json}")

def analyze_ppt(folder: Path):
    """
    对单个学生文件夹中的 PPT 进行配色与排版评估
    结果写入该文件夹下的 PPTScore.txt
    """
    pptx_files = list(folder.glob("*.pptx"))
    if not pptx_files:
        return

    try:
        prs = Presentation(str(pptx_files[0]))
    except Exception as e:
        print(f"无法打开 PPT 文件: {pptx_files[0]}，错误信息: {e}")
        return

    comments = []
    score = 5.0  # 初始分数

    # 1. 模板评估
    try:
        if is_default_office_template(prs):
            score -= 0.5
            comments.append("使用了默认 Office 模板，建议自定义模板提升视觉效果。")
        else:
            comments.append("模板自定义良好，整体风格大方专业。")
    except Exception as e:
        print(f"模板评估出错: {e}")

    # 2. 配色评估
    colors = []
    try:
        fill = prs.slide_master.background.fill
        if fill.type == MSO_FILL_TYPE.SOLID and hasattr(fill.fore_color, 'rgb'):
            rgb = fill.fore_color.rgb
            colors.append((rgb[0], rgb[1], rgb[2]))

        sats = [cv2.cvtColor(np.uint8([[[b, g, r]]]), cv2.COLOR_BGR2HSV)[0, 0, 1] / 255.0 for r, g, b in colors]
        avg_sat = sum(sats) / len(sats) if sats else 0

        if avg_sat < 0.15:
            comments.append("配色柔和，视觉效果舒适。")
        elif avg_sat < 0.35:
            score -= 0.5
            comments.append("配色稍显明亮，建议调整为低饱和度。")
        else:
            score -= 1.0
            comments.append("配色过于鲜艳，建议减少颜色对比度。")
    except Exception as e:
        print(f"配色评估出错: {e}")

    # 3. 形状设计评估
    try:
        total_shapes = shaped_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                total_shapes += 1
                fill = getattr(shape, 'fill', None)
                if fill and fill.type == MSO_FILL_TYPE.GRADIENT:
                    shaped_count += 1

        ratio_shapes = shaped_count / total_shapes if total_shapes else 0
        if ratio_shapes > 0.1:
            comments.append("形状设计丰富，视觉层次感良好。")
        else:
            score -= 0.5
            comments.append("形状设计单调，建议适当添加圆角或渐变元素。")
    except Exception as e:
        print(f"形状设计评估出错: {e}")

    # 4. 文字与图文比例评估
    try:
        line_count = image_area = 0
        slide_area = prs.slide_width * prs.slide_height

        for slide in prs.slides:
            for shp in slide.shapes:
                if hasattr(shp, 'text'):
                    line_count += len(shp.text.strip().splitlines())

            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    image_area += shape.width * shape.height

        avg_lines = line_count / len(prs.slides) if prs.slides else 0
        avg_img_ratio = (image_area / slide_area) / len(prs.slides) if prs.slides else 0

        if avg_lines > 6:
            score -= 0.5
            comments.append(f"文字过多，平均每页 {avg_lines:.1f} 行，建议控制在 6 行以内。")
        else:
            comments.append(f"文字控制良好，平均每页 {avg_lines:.1f} 行。")

        if 0.2 <= avg_img_ratio <= 0.6:
            comments.append(f"图文比例适中，图片覆盖面积：{avg_img_ratio:.2f}。")
        else:
            score -= 0.5
            comments.append(f"图文比例不理想，图片覆盖面积：{avg_img_ratio:.2f}，建议保持在 20%-60%。")
    except Exception as e:
        print(f"文字与图文比例评估出错: {e}")

    # 最终得分
    final_score = max(0, round(score, 1))

    # 写入评分文件
    result_path = folder / "PPTScore.txt"
    with open(result_path, 'w', encoding='utf-8') as f:
        f.write(f"{final_score}/5\n")
        for comment in comments:
            f.write(comment + "\n")


def delete_ppt_files(root_dir):
    """
    删除指定目录下的 .ppt 和 .pptx 文件
    """
    for root, _, files in os.walk(root_dir):
        for file in files:
            if file.lower().endswith(('.ppt', '.pptx')):
                file_path = os.path.join(root, file)
                try:
                    os.remove(file_path)
                    print(f"已删除: {file_path}")
                except Exception as e:
                    print(f"删除失败 {file_path}: {e}")


def a03_02_ppt_scorer(root_dir: str):
    """
    主流程：转换 PPT 到 PPTX，分析并评分，删除原始文件
    """
    convert_ppt_to_pptx(root_dir)

    base_dir = Path(root_dir)
    for student_folder in base_dir.iterdir():
        if student_folder.is_dir():
            analyze_ppt(student_folder)

    # 删除文件
    delete_ppt_files(root_dir)
