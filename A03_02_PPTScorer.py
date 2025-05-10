import os
from pathlib import Path
from pptx import Presentation
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.shapes.picture import Picture
import cv2
import numpy as np

import win32com.client as win32

def convert_ppt_to_pptx(root_dir):
    # 启动 PowerPoint 应用程序
    powerpoint = win32.Dispatch("PowerPoint.Application")
    powerpoint.Visible = False

    for dirpath, _, filenames in os.walk(root_dir):
        for filename in filenames:
            if filename.lower().endswith('.ppt') and not filename.lower().endswith('.pptx'):
                ppt_path = os.path.join(dirpath, filename)
                pptx_path = os.path.join(dirpath, os.path.splitext(filename)[0] + '.pptx')

                print(f"正在转换: {ppt_path} 到 {pptx_path}")

                try:
                    # 打开 .ppt 文件
                    presentation = powerpoint.Presentations.Open(ppt_path)
                    # 转换为 .pptx
                    presentation.SaveAs(pptx_path, 24)  # 24 表示 pptx 格式
                    presentation.Close()

                    # 删除原始 .ppt 文件
                    os.remove(ppt_path)
                    print(f"转换成功并删除: {ppt_path}")

                except Exception as e:
                    print(f"处理文件时发生错误: {ppt_path}\n错误信息: {e}")

    # 退出 PowerPoint 应用程序
    powerpoint.Quit()

def is_default_office_template(prs: Presentation) -> bool:
    """
    判断是否使用了默认 Office 模板
    """
    for master in prs.slide_masters:
        if master.name and "Office Theme" in master.name:
            return True
    return False


def analyze_ppt(folder: Path):
    """
    对单个学生文件夹中的 PPT 进行配色与排版评估
    结果写入该文件夹下的 PPTScore.txt
    """
    pptx_files = list(folder.glob("*.pptx"))
    if not pptx_files:
        return
    prs = Presentation(str(pptx_files[0]))
    comments = []
    score = 5.0  # 起始分为满分

    # 1. 模板评估 (扣 0.5)
    if is_default_office_template(prs):
        score -= 0.5
        comments.append("您使用了默认 Office 模板，建议自定义模板以提升汇报的专业性。")
    else:
        comments.append("模板设计自定义，整体风格大方专业。")

    # 2. 配色评估 (最佳0扣；中等扣0.5；差扣1)
    colors = []
    try:
        fill = prs.slide_master.background.fill
        if fill.type == MSO_FILL_TYPE.SOLID:
            # 只有在 fore_color 有 .rgb 属性时才读取
            if hasattr(fill.fore_color, 'rgb'):
                rgb = fill.fore_color.rgb
                colors.append((rgb[0], rgb[1], rgb[2]))
    except Exception:
        pass

    sats = []
    for r, g, b in colors:
        hsv_val = cv2.cvtColor(
            np.uint8([[[b, g, r]]]),
            cv2.COLOR_BGR2HSV
        )[0, 0, 1] / 255.0
        sats.append(hsv_val)

    avg_sat = sum(sats) / len(sats) if sats else 0
    if avg_sat < 0.15:
        comments.append("配色柔和，契合莫兰迪色系，视觉效果舒适。")
    elif avg_sat < 0.35:
        score -= 0.5
        comments.append("配色较为专业，但色调偏亮，可考虑更多蓝灰元素以提升观感。")
    else:
        score -= 1.0
        comments.append("配色使用较为鲜艳的色彩，建议采用更专业的配色方案，减少视觉干扰。")

    # 3. 形状设计评估 (圆角或渐变至少占比10%得分，否则扣0.5)
    total_shapes = shaped_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            total_shapes += 1
            fill = getattr(shape, 'fill', None)
            if fill and fill.type == MSO_FILL_TYPE.GRADIENT:
                shaped_count += 1
                continue
            try:
                atype = shape.auto_shape_type
                if 'rounded' in atype.name.lower():
                    shaped_count += 1
            except Exception:
                pass

    ratio_shapes = shaped_count / total_shapes if total_shapes else 0
    if ratio_shapes > 0.1:
        comments.append("形状设计较为精致，适当运用了圆角和渐变元素。")
    else:
        score -= 0.5
        comments.append("形状类型较为单一，建议使用圆角、渐变或阴影等设计元素以增加层次感。")

    # 4. 文字与图文比例评估 (文字超6行/页或图文比例不在20%-60%各扣0.5)
    line_count = image_area = 0
    slide_area = prs.slide_width * prs.slide_height
    for slide in prs.slides:
        for shp in slide.shapes:
            if hasattr(shp, 'text') and shp.text.strip():
                line_count += len(shp.text.strip().splitlines())
        for shape in slide.shapes:
            if isinstance(shape, Picture):
                image_area += shape.width * shape.height

    avg_lines = line_count / len(prs.slides) if prs.slides else 0
    avg_img_ratio = (image_area / slide_area) / len(prs.slides) if prs.slides else 0

    if avg_lines <= 6:
        comments.append(f"文字行数控制良好，平均每页{avg_lines:.1f}行，信息简洁明了。")
    else:
        score -= 0.5
        comments.append(f"文字行数偏多，平均每页{avg_lines:.1f}行，建议不超过6行以突出重点。")

    if 0.2 <= avg_img_ratio <= 0.6:
        comments.append(f"图文比例协调，图片覆盖页面面积{avg_img_ratio:.2f}，可视化效果佳。")
    else:
        score -= 0.5
        comments.append(f"图文比例不理想，图片覆盖{avg_img_ratio:.2f}比例，建议保持20%-60%。")

    # 最终得分
    final_score = max(0, round(score, 1))

    # 输出到 PPTScore.txt 位于每个学生文件夹
    result_path = folder / "PPTScore.txt"
    with open(result_path, 'w', encoding='utf-8') as f:
        f.write(f"{final_score}/5\n")
        for sentence in comments:
            f.write(sentence + "\n")


def delete_ppt_files(root_dir):
    """
    删除指定目录下所有 .ppt 和 .pptx 文件
    """
    for root, dirs, files in os.walk(root_dir):
        for file in files:
            if file.lower().endswith(('.ppt', '.pptx')):
                file_path = os.path.join(root, file)
                try:
                    os.remove(file_path)
                except Exception as e:
                    print(f"删除失败 {file_path}: {e}")


def a03_02_ppt_scorer(root_dir: str):
    # root_dir = r"C:\MyPython\ExamScore_AIClass\ExamFiles"
    convert_ppt_to_pptx(root_dir)

    base_dir = Path(root_dir)
    for student_folder in base_dir.iterdir():
        if student_folder.is_dir():
            analyze_ppt(student_folder)

    # 如果不想在评分后删除原始 PPT 文件，可将下面这一行注释掉
    delete_ppt_files(r"C:\MyPython\ExamScore_AIClass\ExamFiles")


